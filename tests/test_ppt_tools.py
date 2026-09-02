"""PPT 渲染工具桥测试（backend_ppt_tools）。

渲染测试依赖本机 LibreOffice / PowerPoint COM，缺失时自动跳过。
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient

import backend_ppt_tools as ppt_tools


def _make_client() -> TestClient:
    app = FastAPI()
    app.include_router(ppt_tools.router)
    return TestClient(app)


def _make_sample_pptx() -> bytes:
    pptx = pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "ZonKey 渲染测试"
    slide.placeholders[1].text = "sample subtitle"

    bullet_layout = prs.slide_layouts[1]
    second = prs.slides.add_slide(bullet_layout)
    body = second.placeholders[1].text_frame
    body.text = "第一行要点"
    para = body.add_paragraph()
    para.text = "第二行要点"
    para.font.size = Pt(18)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


@pytest.fixture(scope="module")
def sample_pptx() -> bytes:
    return _make_sample_pptx()


@pytest.fixture(scope="module")
def has_renderer() -> bool:
    capability = ppt_tools.render_capability()
    return bool(capability["libreoffice"] or capability["powerpoint_com"])


def test_capability_flags():
    data = ppt_tools.render_capability()
    assert isinstance(data["libreoffice"], bool)
    assert isinstance(data["powerpoint_com"], bool)


def test_render_rejects_non_pptx():
    client = _make_client()
    response = client.post(
        "/api/ppt/render",
        files={"file": ("notes.txt", io.BytesIO(b"not a pptx"), "text/plain")},
        data={"target": "pdf"},
    )
    assert response.status_code == 400


def test_render_rejects_fake_zip():
    client = _make_client()
    response = client.post(
        "/api/ppt/render",
        files={"file": ("fake.pptx", io.BytesIO(b"PK\x03\x04garbage"), "application/octet-stream")},
        data={"target": "pdf"},
    )
    assert response.status_code in {400, 503}


@pytest.mark.parametrize("target", ["pdf", "images"])
def test_render_end_to_end(target: str, sample_pptx: bytes, has_renderer: bool, tmp_path: Path):
    if not has_renderer:
        pytest.skip("本机无 LibreOffice / PowerPoint 渲染器")

    client = _make_client()
    response = client.post(
        "/api/ppt/render",
        files={"file": ("zonkey_test_deck.pptx", io.BytesIO(sample_pptx), "application/octet-stream")},
        data={"target": target, "dpi": "96"},
    )
    assert response.status_code == 200, response.text
    payload = response.json()
    assert payload["status"] == "success"
    assert payload["renderer"] in {"libreoffice", "powerpoint-com"}

    output_path = Path(payload["output_dir"]) / payload["download_name"]
    # 测试产物立即校验后清理，不污染 output/
    try:
        assert output_path.exists(), f"产物缺失: {output_path}"
        if target == "pdf":
            from core.pdfio import page_count

            assert page_count(str(output_path)) >= 2
        else:
            with zipfile.ZipFile(output_path) as bundle:
                names = [name for name in bundle.namelist() if name.endswith(".png")]
                assert len(names) >= 2
    finally:
        output_path.unlink(missing_ok=True)
