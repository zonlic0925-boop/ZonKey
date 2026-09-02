"""为全功能 UI 实测生成合成样本（多页 PDF / PPTX / 图片 / 音频 / 视频）。"""
import os
import struct
import wave
import math

HERE = os.path.dirname(os.path.abspath(__file__))

# ---------- 多页 PDF（merge/split/extract/rotate/page-numbers/organize/editor） ----------
from reportlab.pdfgen import canvas

mp = os.path.join(HERE, "multi_page.pdf")
c = canvas.Canvas(mp, pagesize=(595, 842))
for i in range(1, 5):
    c.setFont("Helvetica", 24)
    c.drawString(80, 700, f"Multi Page Document - Page {i}")
    c.setFont("Helvetica", 12)
    for ln in range(6):
        c.drawString(80, 650 - ln * 24, f"Page {i} line {ln + 1}: sample body text for tool testing.")
    c.showPage()
c.save()
print("multi_page.pdf OK")

# ---------- 表格 PDF（pdf-to-excel） ----------
tp = os.path.join(HERE, "sample_table.pdf")
c = canvas.Canvas(tp, pagesize=(595, 842))
rows = [
    ["Item", "Qty", "Price"],
    ["Widget", "10", "9.99"],
    ["Gadget", "3", "24.50"],
    ["Gizmo", "7", "13.20"],
]
y = 780
for r in rows:
    x = 80
    for cell in r:
        c.drawString(x, y, cell)
        x += 120
    y -= 30
c.save()
print("sample_table.pdf OK")

# ---------- PPTX（PPT 工坊） ----------
from pptx import Presentation
from pptx.util import Inches, Pt

pp = os.path.join(HERE, "sample.pptx")
prs = Presentation()
layout = prs.slide_layouts[0]
s1 = prs.slides.add_slide(layout)
s1.shapes.title.text = "ZonKey PPT Test"
s1.placeholders[1].text = "Automated UI test deck"
layout2 = prs.slide_layouts[1]
for i in range(2, 4):
    s = prs.slides.add_slide(layout2)
    s.shapes.title.text = f"Slide {i}"
    body = s.placeholders[1].text_frame
    body.text = f"Bullet one for slide {i}"
    p = body.add_paragraph()
    p.text = f"Bullet two for slide {i}"
prs.save(pp)
print("sample.pptx OK")

# ---------- 图片（图像工坊） ----------
from PIL import Image, ImageDraw

img = Image.new("RGB", (640, 480), (240, 90, 90))
d = ImageDraw.Draw(img)
d.rectangle([80, 80, 560, 400], fill=(90, 140, 240))
d.ellipse([240, 160, 400, 320], fill=(250, 220, 80))
img.save(os.path.join(HERE, "sample_image.png"))
img.save(os.path.join(HERE, "sample_image.jpg"), quality=90)

img2 = Image.new("RGB", (400, 300), (120, 200, 120))
ImageDraw.Draw(img2).text((30, 30), "Second image", fill=(20, 40, 20))
img2.save(os.path.join(HERE, "sample_image2.png"))
print("images OK")

# ---------- 图标源图（icon-gen） ----------
ic = Image.new("RGBA", (256, 256), (0, 0, 0, 0))
di = ImageDraw.Draw(ic)
di.rounded_rectangle([16, 16, 240, 240], radius=48, fill=(40, 90, 200, 255))
di.ellipse([88, 88, 168, 168], fill=(255, 255, 255, 255))
ic.save(os.path.join(HERE, "sample_icon.png"))
print("sample_icon.png OK")

# ---------- 音频（音视频中心：bpm/clip/convert/extract） ----------
wav_path = os.path.join(HERE, "sample_audio.wav")
sr = 22050
with wave.open(wav_path, "w") as w:
    w.setnchannels(1)
    w.setsampwidth(2)
    w.setframerate(sr)
    frames = bytearray()
    # 4 秒：1 秒 440Hz + 1 秒静音 交替（可测 BPM 与裁剪）
    for t in range(sr * 4):
        beat = (t // (sr // 2)) % 2  # 每半秒交替
        v = math.sin(2 * math.pi * 440 * (t % (sr // 2)) / sr) * (12000 if beat else 0)
        frames += struct.pack("<h", int(v))
    w.writeframes(bytes(frames))
print("sample_audio.wav OK")

# ---------- 视频（video-convert / frame / gif） ----------
mp4 = os.path.join(HERE, "sample_video.mp4")
os.system(
    f'ffmpeg -y -loglevel error -f lavfi -i "testsrc=duration=3:size=320x240:rate=15" '
    f'-f lavfi -i "sine=frequency=440:duration=3" -pix_fmt yuv420p -shortest "{mp4}"'
)
print("sample_video.mp4 OK" if os.path.exists(mp4) else "sample_video.mp4 FAILED")
