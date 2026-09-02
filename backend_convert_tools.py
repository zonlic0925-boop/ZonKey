"""文档转换工具桥（转换引擎 P2 垂直切片，宽松许可选型）。

- PDF→Word：pdfplumber（文本行+表格+坐标）→ python-docx 重建（自研链路，
  刻意规避 pdf2docx 的 PyMuPDF 硬依赖；复杂版式保真度下降，UI 需如实标注）。
- PDF→Excel：pdfplumber 表格抽取 → openpyxl（每页一个工作表）。
- PDF→PPT：pypdfium2 逐页渲染 → python-pptx 整页贴图（视觉版式还原，
  非可编辑文本还原）。
- Office→PDF：Word/Excel COM 导出（保真优先）；无 Office 时任务报错明示，
  能力端点只做导入级探测，禁止静默降级（mammoth/reportlab HTML 兜底链后续接入）。
- PDF 修复：pikepdf（qpdf 内核）损坏结构恢复重写。

全模块不依赖 PyMuPDF（AGPL 退出计划，见 docs/ILOVEPDF_INTEGRATION_PLAN.md Phase M）。

长任务采用 job 模式：POST 立即返回 job_id，后台线程执行，
前端轮询 GET /api/convert/jobs/{job_id} 取进度与产物。
产物写入 output/，经 /api/download/{filename} 与原生另存为取件，
不经浏览器 blob 下载通道（pywebview 壳兼容）。
"""

from __future__ import annotations

import re
import shutil
import statistics
import tempfile
import threading
import uuid
from html.parser import HTMLParser
from importlib.metadata import version as _dist_version
from pathlib import Path
from typing import Any, Optional

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from reportlab.platypus import Paragraph

from core.app_paths import ensure_runtime_layout

router = APIRouter(prefix="/api/convert", tags=["convert-tools"])

PROJECT_ROOT = ensure_runtime_layout()
OUTPUT_DIR = PROJECT_ROOT / "output"

MAX_INPUT_BYTES = 200 * 1024 * 1024
MAX_PAGES = 300
PDF_DPI_RANGE = (72, 300)

_WORD_EXTS = {".docx", ".doc"}
_EXCEL_EXTS = {".xlsx", ".xls"}

_jobs: dict[str, dict[str, Any]] = {}
_jobs_lock = threading.Lock()
# Word/Excel/PowerPoint COM 共用一把锁：避免多任务同时拉起 Office 实例互相拖垮
_com_lock = threading.Lock()


# ---------------------------------------------------------------------------
# 基础设施：上传落盘 / 产物命名 / 任务登记
# ---------------------------------------------------------------------------

def _safe_base_name(source_name: str) -> str:
    base = Path(source_name).stem.strip()
    cleaned = "".join("_" if ch in '\\/:*?"<>|' else ch for ch in base).strip("._ ") or "document"
    return cleaned[:80]


def _unique_output_path(base: str, suffix: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    candidate = OUTPUT_DIR / f"{base}{suffix}"
    if not candidate.exists():
        return candidate
    for index in range(1, 1000):
        candidate = OUTPUT_DIR / f"{base}_{index}{suffix}"
        if not candidate.exists():
            return candidate
    return OUTPUT_DIR / f"{base}_{uuid.uuid4().hex[:6]}{suffix}"


def _write_upload(tmp_dir: Path, file: UploadFile, raw_name: str) -> Path:
    dest = tmp_dir / f"src_{uuid.uuid4().hex[:8]}_{raw_name}"
    try:
        with open(dest, "wb") as handle:
            shutil.copyfileobj(file.file, handle)
    except OSError as exc:
        raise HTTPException(status_code=400, detail=f"上传文件写入失败: {exc}") from exc
    if dest.stat().st_size > MAX_INPUT_BYTES:
        raise HTTPException(status_code=400, detail="文件超过 200MB 上限")
    return dest


def _save_pdf_upload(tmp_dir: Path, file: UploadFile) -> tuple[Path, str]:
    raw_name = Path(file.filename or "document.pdf").name
    if not raw_name.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="仅支持 .pdf 文件")
    dest = _write_upload(tmp_dir, file, raw_name)
    if dest.stat().st_size < 100:
        raise HTTPException(status_code=400, detail="PDF 文件无效（体积过小）")
    with open(dest, "rb") as handle:
        if b"%PDF-" not in handle.read(1024):
            raise HTTPException(status_code=400, detail="PDF 文件无效（缺少 %PDF 头）")
    return dest, _safe_base_name(raw_name)


def _save_office_upload(
    tmp_dir: Path, file: UploadFile, allowed: str | None = None
) -> tuple[Path, str, str]:
    raw_name = Path(file.filename or "document.docx").name
    ext = Path(raw_name).suffix.lower()
    if ext in _WORD_EXTS:
        kind = "word"
    elif ext in _EXCEL_EXTS:
        kind = "excel"
    elif ext in {".pptx", ".ppt"}:
        raise HTTPException(status_code=400, detail="PPT 转 PDF 请使用 PPT 工坊（已支持 LibreOffice/COM 双路渲染）")
    else:
        raise HTTPException(status_code=400, detail=f"不支持的格式: {ext or '(无扩展名)'}，支持 .docx/.doc/.xlsx/.xls")
    if allowed is not None and kind != allowed:
        raise HTTPException(
            status_code=400,
            detail=f"文件类型不匹配：本工具接受 {('Word (.docx/.doc)' if allowed == 'word' else 'Excel (.xlsx/.xls)')}",
        )
    dest = _write_upload(tmp_dir, file, raw_name)
    return dest, _safe_base_name(raw_name), kind


def _set_job(job_id: str, **fields: Any) -> None:
    with _jobs_lock:
        _jobs[job_id].update(fields)


def _get_job(job_id: str) -> Optional[dict[str, Any]]:
    with _jobs_lock:
        return dict(_jobs.get(job_id) or {})


def _prune_jobs() -> None:
    with _jobs_lock:
        if len(_jobs) > 64:
            done_ids = [k for k, v in _jobs.items() if v.get("status") in ("done", "error")]
            for k in done_ids[:-32]:
                _jobs.pop(k, None)


def _new_job(kind: str) -> str:
    job_id = uuid.uuid4().hex[:12]
    with _jobs_lock:
        _jobs[job_id] = {
            "job_id": job_id, "kind": kind, "status": "running",
            "progress": 5, "stage": "prepare", "error": None, "outputs": [],
        }
    _prune_jobs()
    return job_id


def _spawn_job(job_id: str, tmp_dir: Path, target: Any) -> None:
    def worker() -> None:
        try:
            target()
        except Exception as exc:  # noqa: BLE001 — 任务线程内统一兜底为 error 状态
            _set_job(job_id, status="error", error=str(exc))
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    threading.Thread(target=worker, daemon=True).start()


# ---------------------------------------------------------------------------
# 能力探测
# ---------------------------------------------------------------------------

def _probe_win32() -> bool:
    try:
        import pythoncom  # noqa: F401
        import win32com.client  # noqa: F401
        return True
    except ImportError:
        return False


def _engine_version(dist_name: str) -> Optional[str]:
    try:
        return _dist_version(dist_name)
    except Exception:
        return None


@router.get("/capability")
def convert_capability() -> dict[str, Any]:
    win32_ready = _probe_win32()
    rapidocr_ready = False
    try:
        from core.detector.ocr_channel import get_ocr_engine  # noqa: F401

        rapidocr_ready = True
    except Exception:  # noqa: BLE001
        rapidocr_ready = False
    return {
        "engines": {
            name: _engine_version(dist)
            for name, dist in {
                "pdfplumber": "pdfplumber",
                "python-docx": "python-docx",
                "openpyxl": "openpyxl",
                "pypdfium2": "pypdfium2",
                "python-pptx": "python-pptx",
                "pikepdf": "pikepdf",
                "reportlab": "reportlab",
                "mammoth": "mammoth",
            }.items()
        },
        # 导入级探测（与 PPT 工坊 /api/ppt/render/capability 同口径）；
        # 真实可用性在转换时确认，缺 Office 组件时任务报错明示
        "pywin32": win32_ready,
        "word_com": win32_ready,
        "excel_com": win32_ready,
        "rapidocr": rapidocr_ready,
        "word_fallback": _engine_version("mammoth") is not None,
    }


@router.get("/jobs/{job_id}")
def convert_job(job_id: str) -> dict[str, Any]:
    job = _get_job(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="任务不存在或已过期")
    return job


# ---------------------------------------------------------------------------
# PDF→Word（pdfplumber 文本行+表格 → python-docx 重建）
# ---------------------------------------------------------------------------

def _is_cjk_char(ch: str) -> bool:
    return bool(ch) and "\u4e00" <= ch <= "\u9fff"


def _heading_level(size: float, body_size: float) -> int:
    if body_size <= 0:
        return 0
    ratio = size / body_size
    if ratio >= 1.6:
        return 1
    if ratio >= 1.3:
        return 2
    if ratio >= 1.15:
        return 3
    return 0


def _page_text_lines(page: Any) -> tuple[list[dict[str, Any]], list[Any]]:
    """按行抽取文本行（字号/粗体/坐标），剔除落在表格框内的行；返回 (lines, tables)。"""
    tables = page.find_tables()
    boxes = [t.bbox for t in tables if t.bbox]
    try:
        raw_lines = page.extract_text_lines()
    except Exception:
        raw_lines = []
    lines: list[dict[str, Any]] = []
    for line in raw_lines:
        text = (line.get("text") or "").strip()
        if not text:
            continue
        cx = (line["x0"] + line["x1"]) / 2
        cy = (line["top"] + line["bottom"]) / 2
        if any(x0 <= cx <= x1 and top <= cy <= bottom for x0, top, x1, bottom in boxes):
            continue
        chars = [c for c in (line.get("chars") or []) if (c.get("text") or "").strip()]
        sizes = [c.get("size", 10.0) for c in chars]
        fonts = [c.get("fontname", "") for c in chars]
        line["_size"] = max(sizes) if sizes else 10.0
        line["_bold"] = bool(fonts) and sum(1 for f in fonts if "bold" in f.lower()) >= len(fonts) * 0.6
        lines.append(line)
    lines.sort(key=lambda l: (round(l["top"], 1), l["x0"]))
    return lines, tables


def _lines_to_paragraphs(lines: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """行→段落：同级行按行距与左对齐归并，标题行独立成段。"""
    body_size = statistics.median([l["_size"] for l in lines]) if lines else 0.0
    paragraphs: list[dict[str, Any]] = []
    current: Optional[dict[str, Any]] = None
    for line in lines:
        level = _heading_level(line["_size"], body_size)
        height = max(line["bottom"] - line["top"], 1.0)
        text = (line.get("text") or "").strip()
        if (
            current is not None
            and level == 0 == current["level"]
            and line["top"] - current["bottom"] <= 0.6 * current["height"]
            and abs(line["x0"] - current["x0"]) <= 3
        ):
            joiner = "" if _is_cjk_char(current["text"][-1:]) and _is_cjk_char(text[:1]) else " "
            current["text"] += joiner + text
            current["bottom"] = line["bottom"]
            current["height"] = max(current["height"], height)
        else:
            current = {
                "text": text, "level": level, "bold": line["_bold"],
                "x0": line["x0"], "top": line["top"],
                "bottom": line["bottom"], "height": height,
            }
            paragraphs.append(current)
    return paragraphs


def _normalize_table_rows(raw_rows: list[Any]) -> list[list[str]]:
    width = max((len(r) for r in raw_rows if r), default=0)
    rows: list[list[str]] = []
    for raw in raw_rows:
        row = [str(cell).replace("\n", " ").strip() if cell is not None else "" for cell in (raw or [])]
        row += [""] * (width - len(row))
        rows.append(row)
    return rows


def _add_docx_paragraph(doc: Any, para: dict[str, Any]) -> None:
    text = para["text"].strip()
    if not text:
        return
    if para["level"] > 0:
        doc.add_heading(text, level=para["level"])
        return
    run = doc.add_paragraph().add_run(text)
    if para["bold"]:
        run.bold = True


def _add_docx_table(doc: Any, rows: list[list[str]]) -> None:
    cols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=cols)
    try:
        table.style = "Table Grid"
    except Exception:
        pass  # 模板样式缺失不影响数据写入
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r == 0 and val and cell.paragraphs[0].runs:
                cell.paragraphs[0].runs[0].bold = True
    doc.add_paragraph()  # 表后空行，避免紧贴后续段落


def _build_docx_from_pdf(source: Path, job_id: str) -> tuple[Any, dict[str, int], list[str]]:
    import pdfplumber
    from docx import Document

    doc = Document()
    stats = {"pages": 0, "tables": 0, "paragraphs": 0}
    warnings: list[str] = []

    with pdfplumber.open(str(source)) as pdf:
        total = len(pdf.pages)
        if total == 0:
            raise RuntimeError("PDF 无页面")
        if total > MAX_PAGES:
            raise RuntimeError(f"PDF 超过 {MAX_PAGES} 页上限")
        for index, page in enumerate(pdf.pages, start=1):
            _set_job(job_id, progress=10 + int(80 * index / total), stage=f"converting {index}/{total}")
            lines, tables = _page_text_lines(page)
            if not lines and not tables:
                warnings.append(f"第 {index} 页未提取到文本与表格（可能为扫描页）")
            paragraphs = _lines_to_paragraphs(lines)
            elements = [("para", p["top"], p) for p in paragraphs]
            elements += [("table", t.bbox[1], t) for t in tables]
            for kind, _, payload in sorted(elements, key=lambda e: e[1]):
                if kind == "para":
                    _add_docx_paragraph(doc, payload)
                    stats["paragraphs"] += 1
                else:
                    rows = _normalize_table_rows(payload.extract() or [])
                    if rows and rows[0]:
                        _add_docx_table(doc, rows)
                        stats["tables"] += 1
                    else:
                        warnings.append(f"第 {index} 页检测到空表格，已跳过")
            if index < total:
                doc.add_page_break()
            stats["pages"] = index

    if stats["paragraphs"] == 0 and stats["tables"] == 0:
        raise RuntimeError("PDF 无文本层（可能是扫描件）——请改用「OCR 导出」工具提取文字")
    return doc, stats, warnings


@router.post("/pdf-to-word")
def pdf_to_word(file: UploadFile = File(...)) -> dict[str, Any]:
    tmp_dir = Path(tempfile.mkdtemp(prefix="convert_word_"))
    source, base = _save_pdf_upload(tmp_dir, file)
    out_path = _unique_output_path(base, ".docx")
    job_id = _new_job("pdf-to-word")

    def target() -> None:
        doc, stats, warnings = _build_docx_from_pdf(source, job_id)
        doc.save(str(out_path))
        _set_job(
            job_id, status="done", progress=100, stage="done",
            engine="pdfplumber+python-docx",
            page_count=stats["pages"], table_count=stats["tables"],
            paragraph_count=stats["paragraphs"], warnings=warnings,
            outputs=[{"name": out_path.name, "dir": str(out_path.parent)}],
        )

    _spawn_job(job_id, tmp_dir, target)
    return {"job_id": job_id, "status": "running"}


# ---------------------------------------------------------------------------
# PDF→Excel（pdfplumber 表格抽取 → openpyxl）
# ---------------------------------------------------------------------------

def _display_width(text: str) -> int:
    return sum(2 if ord(ch) > 0x2E80 else 1 for ch in text)


def _smart_value(text: str) -> Any:
    """纯数字单元格转数值类型，其余保持字符串（保留前导零等原貌）。"""
    if re.fullmatch(r"-?\d+", text):
        try:
            return int(text)
        except ValueError:
            return text
    if re.fullmatch(r"-?\d+\.\d+", text):
        try:
            return float(text)
        except ValueError:
            return text
    return text


def _build_xlsx_from_pdf(source: Path, out_path: Path, job_id: str) -> tuple[int, int, list[str]]:
    import pdfplumber
    from openpyxl import Workbook
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    workbook = Workbook()
    workbook.remove(workbook.active)
    table_count = 0
    page_count = 0
    warnings: list[str] = []

    with pdfplumber.open(str(source)) as pdf:
        total = len(pdf.pages)
        if total == 0:
            raise RuntimeError("PDF 无页面")
        if total > MAX_PAGES:
            raise RuntimeError(f"PDF 超过 {MAX_PAGES} 页上限")
        for index, page in enumerate(pdf.pages, start=1):
            _set_job(job_id, progress=10 + int(80 * index / total), stage=f"extracting {index}/{total}")
            try:
                tables = page.find_tables()
            except Exception:
                tables = []
            if not tables:
                continue
            sheet = workbook.create_sheet(f"Page{index}")
            cursor = 1
            widths: dict[int, int] = {}
            for table in tables:
                rows = _normalize_table_rows(table.extract() or [])
                if not rows or not any(any(cell for cell in row) for row in rows):
                    continue
                for r_idx, row in enumerate(rows, start=1):
                    for c_idx, val in enumerate(row, start=1):
                        cell = sheet.cell(row=cursor + r_idx - 1, column=c_idx, value=_smart_value(val))
                        if r_idx == 1 and len(rows) > 1:
                            cell.font = Font(bold=True)
                        widths[c_idx] = max(widths.get(c_idx, 0), _display_width(val))
                cursor += len(rows) + 1
                table_count += 1
            for c_idx, width in widths.items():
                sheet.column_dimensions[get_column_letter(c_idx)].width = min(60, max(10, width + 2))
            page_count = index

    if table_count == 0:
        raise RuntimeError("未检测到表格——本工具仅提取带边框/线条的表格，正文内容请使用 PDF→Word")
    workbook.save(str(out_path))
    return table_count, page_count, warnings


@router.post("/pdf-to-excel")
def pdf_to_excel(file: UploadFile = File(...)) -> dict[str, Any]:
    tmp_dir = Path(tempfile.mkdtemp(prefix="convert_excel_"))
    source, base = _save_pdf_upload(tmp_dir, file)
    out_path = _unique_output_path(f"{base}_tables", ".xlsx")
    job_id = _new_job("pdf-to-excel")

    def target() -> None:
        table_count, page_count, warnings = _build_xlsx_from_pdf(source, out_path, job_id)
        _set_job(
            job_id, status="done", progress=100, stage="done",
            engine="pdfplumber+openpyxl", page_count=page_count, table_count=table_count,
            warnings=warnings,
            outputs=[{"name": out_path.name, "dir": str(out_path.parent)}],
        )

    _spawn_job(job_id, tmp_dir, target)
    return {"job_id": job_id, "status": "running"}


# ---------------------------------------------------------------------------
# PDF→PPT（pypdfium2 逐页渲染 → python-pptx 整页贴图）
# ---------------------------------------------------------------------------

@router.post("/pdf-to-ppt")
def pdf_to_ppt(
    file: UploadFile = File(...),
    dpi: int = Form(150),
    image_format: str = Form("png"),
) -> dict[str, Any]:
    dpi_clamped = min(max(int(dpi), PDF_DPI_RANGE[0]), PDF_DPI_RANGE[1])
    fmt = "jpeg" if image_format == "jpeg" else "png"

    tmp_dir = Path(tempfile.mkdtemp(prefix="convert_ppt_"))
    source, base = _save_pdf_upload(tmp_dir, file)
    out_path = _unique_output_path(f"{base}_slides", ".pptx")
    job_id = _new_job("pdf-to-ppt")

    def target() -> None:
        import pypdfium2 as pdfium
        from pptx import Presentation
        from pptx.util import Emu

        source_pdf = pdfium.PdfDocument(str(source))
        try:
            total = len(source_pdf)
            if total == 0:
                raise RuntimeError("PDF 无页面")
            if total > MAX_PAGES:
                raise RuntimeError(f"PDF 超过 {MAX_PAGES} 页上限")
            width_pt, height_pt = source_pdf[0].get_size()
            prs = Presentation()
            prs.slide_width = Emu(int(round(width_pt * 12700)))
            prs.slide_height = Emu(int(round(height_pt * 12700)))
            blank_layout = prs.slide_layouts[6]

            ext = "jpg" if fmt == "jpeg" else "png"
            for index in range(total):
                _set_job(job_id, progress=10 + int(80 * (index + 1) / total), stage=f"rendering {index + 1}/{total}")
                bitmap = source_pdf[index].render(scale=dpi_clamped / 72.0)
                pil_image = bitmap.to_pil()
                image_path = tmp_dir / f"page_{index + 1}.{ext}"
                if ext == "jpg":
                    pil_image.convert("RGB").save(image_path, "JPEG", quality=90)
                else:
                    pil_image.save(image_path, "PNG")
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(image_path), 0, 0,
                    width=prs.slide_width, height=prs.slide_height,
                )
            prs.save(str(out_path))
        finally:
            source_pdf.close()

        _set_job(
            job_id, status="done", progress=100, stage="done",
            engine="pypdfium2+python-pptx", page_count=total,
            dpi=dpi_clamped, image_format=fmt,
            note="视觉版式还原（整页贴图），不含可编辑文本",
            outputs=[{"name": out_path.name, "dir": str(out_path.parent)}],
        )

    _spawn_job(job_id, tmp_dir, target)
    return {"job_id": job_id, "status": "running"}


# ---------------------------------------------------------------------------
# Office→PDF（Word/Excel COM 导出，保真优先）
# ---------------------------------------------------------------------------

def _convert_office_with_com(source: Path, out_path: Path, kind: str, job_id: str) -> None:
    if not _probe_win32():
        raise RuntimeError("缺少 pywin32，无法调用 Office COM（请安装 requirements 中的 pywin32）")
    import pythoncom
    import win32com.client

    with _com_lock:
        pythoncom.CoInitialize()
        app = None
        handle = None
        try:
            if kind == "word":
                app = win32com.client.DispatchEx("Word.Application")
                app.Visible = False
                app.DisplayAlerts = 0
                _set_job(job_id, progress=30, stage="word-export")
                handle = app.Documents.Open(str(source), False, True, False)  # ConfirmConversions, ReadOnly, AddToRecentFiles
                if handle is None:
                    raise RuntimeError("Word 无法打开该文档")
                handle.SaveAs2(str(out_path), 17)  # wdFormatPDF
            else:
                app = win32com.client.DispatchEx("Excel.Application")
                app.Visible = False
                app.DisplayAlerts = False
                _set_job(job_id, progress=30, stage="excel-export")
                handle = app.Workbooks.Open(str(source), 0, True)  # UpdateLinks=0, ReadOnly
                if handle is None:
                    raise RuntimeError("Excel 无法打开该工作簿")
                handle.ExportAsFixedFormat(0, str(out_path))  # xlTypePDF
            _set_job(job_id, progress=90, stage="finalize")
        finally:
            try:
                if handle is not None:
                    handle.Close(False)
            except Exception:
                pass
            try:
                if app is not None:
                    app.Quit()
            except Exception:
                pass
            pythoncom.CoUninitialize()


def _word_to_pdf_fallback(source: Path, out_path: Path, job_id: str) -> None:
    """Word 兜底：mammoth docx→HTML → platypus（保真度下降，正文/标题/表格保留）。"""
    import mammoth

    _set_job(job_id, progress=30, stage="fallback-html")
    with open(source, "rb") as docx_stream:
        result = mammoth.convert_to_html(docx_stream)
    warnings = list(result.messages and ["忽略: " + m.message for m in result.messages[:5]] or [])
    flowables, extra = _html_to_flowables(result.value, _ensure_pdf_font())
    warnings.extend(extra)
    _set_job(job_id, progress=70, stage="fallback-pdf")
    _build_pdf_from_flowables(flowables, out_path, source.stem)


def _excel_to_pdf_fallback(source: Path, out_path: Path, job_id: str) -> None:
    """Excel 兜底：openpyxl 读单元格 → HTML 表格 → platypus（公式取缓存值，样式不保留）。"""
    from openpyxl import load_workbook

    _set_job(job_id, progress=30, stage="fallback-read")
    workbook = load_workbook(str(source), data_only=True, read_only=True)
    sections: list[str] = []
    for sheet in workbook.worksheets:
        rows_html: list[str] = []
        for row in sheet.iter_rows(max_row=500, max_col=30, values_only=True):
            if not any(v is not None and str(v).strip() for v in row):
                continue
            cells = "".join(f"<td>{_esc(str(v))}</td>" for v in row)
            rows_html.append(f"<tr>{cells}</tr>")
        if rows_html:
            sections.append(f"<h2>{_esc(sheet.title)}</h2><table>{''.join(rows_html)}</table>")
    workbook.close()
    if not sections:
        raise RuntimeError("工作簿无可见数据")
    flowables, _extra = _html_to_flowables("".join(sections), _ensure_pdf_font())
    _set_job(job_id, progress=70, stage="fallback-pdf")
    _build_pdf_from_flowables(flowables, out_path, source.stem)


def _start_office_job(file: UploadFile, allowed: str | None, op_name: str) -> dict[str, Any]:
    tmp_dir = Path(tempfile.mkdtemp(prefix="convert_office_"))
    source, base, kind = _save_office_upload(tmp_dir, file, allowed)
    out_path = _unique_output_path(base, ".pdf")
    job_id = _new_job(op_name)

    def target() -> None:
        try:
            _convert_office_with_com(source, out_path, kind, job_id)
            engine = "office-com"
            fallback_note = None
        except Exception as com_exc:
            # 兜底链（诚实标注保真度下降，禁止静默降级）：Word → mammoth HTML →
            # reportlab；Excel → openpyxl → HTML 表格 → reportlab
            try:
                if kind == "word":
                    _word_to_pdf_fallback(source, out_path, job_id)
                    engine = "mammoth+reportlab"
                else:
                    _excel_to_pdf_fallback(source, out_path, job_id)
                    engine = "openpyxl+reportlab"
            except Exception as fb_exc:
                raise RuntimeError(
                    f"Office COM 导出失败且兜底转换失败。"
                    f"COM（需本机 Microsoft {('Word' if kind == 'word' else 'Excel')}）: {com_exc}；"
                    f"兜底: {fb_exc}"
                ) from fb_exc
            fallback_note = (
                "本机 Office COM 不可用，已用 HTML 兜底转换：正文与表格保留，"
                "复杂排版/图片/公式保真度下降"
            )
        if not out_path.exists() or out_path.stat().st_size == 0:
            raise RuntimeError("转换未产出有效 PDF")
        payload: dict[str, Any] = {
            "status": "done", "progress": 100, "stage": "done",
            "engine": engine,
            "office_app": ("word" if kind == "word" else "excel"),
            "outputs": [{"name": out_path.name, "dir": str(out_path.parent)}],
        }
        if fallback_note:
            payload["note"] = fallback_note
            payload["warnings"] = [fallback_note]
        _set_job(job_id, **payload)

    _spawn_job(job_id, tmp_dir, target)
    return {"job_id": job_id, "status": "running"}


@router.post("/office-to-pdf")
def office_to_pdf(file: UploadFile = File(...)) -> dict[str, Any]:
    """兼容端点：Word/Excel 混合上传（保留旧接口，前端已改用拆分端点）。"""
    return _start_office_job(file, None, "office-to-pdf")


@router.post("/word-to-pdf")
def word_to_pdf(file: UploadFile = File(...)) -> dict[str, Any]:
    """Word → PDF（.docx/.doc）：Office COM 优先，mammoth+reportlab 兜底。"""
    return _start_office_job(file, "word", "word-to-pdf")


@router.post("/excel-to-pdf")
def excel_to_pdf(file: UploadFile = File(...)) -> dict[str, Any]:
    """Excel → PDF（.xlsx/.xls）：Office COM 优先，openpyxl+reportlab 兜底。"""
    return _start_office_job(file, "excel", "excel-to-pdf")


# ---------------------------------------------------------------------------
# PDF 修复（pikepdf / qpdf 损坏恢复）
# ---------------------------------------------------------------------------

@router.post("/repair")
def repair_pdf(file: UploadFile = File(...)) -> dict[str, Any]:
    try:
        import pikepdf
    except ImportError as exc:
        raise HTTPException(status_code=503, detail="修复引擎未安装：请执行 pip install pikepdf 后重启应用。") from exc

    with tempfile.TemporaryDirectory(prefix="convert_repair_") as tmp:
        tmp_dir = Path(tmp)
        source, base = _save_pdf_upload(tmp_dir, file)
        out_path = _unique_output_path(f"{base}_repaired", ".pdf")
        try:
            with pikepdf.open(str(source)) as pdf:  # 打开时自动尝试恢复损坏的 xref/对象结构
                page_count = len(pdf.pages)
                pdf.save(str(out_path))
        except Exception as exc:
            raise HTTPException(status_code=422, detail=f"PDF 损坏严重，修复失败: {exc}") from exc
        if not out_path.exists() or out_path.stat().st_size == 0:
            raise HTTPException(status_code=500, detail="修复产物为空")
        return {
            "status": "success",
            "engine": "pikepdf (qpdf)",
            "page_count": page_count,
            "download_name": out_path.name,
            "output_dir": str(out_path.parent),
        }


# ---------------------------------------------------------------------------
# 深度压缩（栅格化重编码：pypdfium2 渲染 + PIL 量化降色 + JPEG DCTDecode
# 直嵌 pikepdf，不经过 reportlab 无损重包——reportlab 会把 JPEG 解码成裸
# 像素再 Flate 回写，压缩因子完全丢失，文字型 PDF 反而变大。
# 若压缩后体积更大则退回原文件（保底返原）。
# ---------------------------------------------------------------------------

@router.post("/compress-deep")
def compress_deep(
    file: UploadFile = File(...),
    dpi: int = Form(144),
    quality: int = Form(70),
) -> dict[str, Any]:
    dpi_clamped = min(max(int(dpi), 72), 200)
    quality_clamped = min(max(int(quality), 30), 95)

    tmp_dir = Path(tempfile.mkdtemp(prefix="convert_compress_"))
    source, base = _save_pdf_upload(tmp_dir, file)
    out_path = _unique_output_path(f"{base}_compressed", ".pdf")
    original_bytes = source.stat().st_size
    job_id = _new_job("compress-deep")

    def target() -> None:
        import io as _io

        import pikepdf
        import pypdfium2 as pdfium
        from PIL import Image
        from pikepdf import Array, Dictionary, Name

        source_pdf = pdfium.PdfDocument(str(source))
        try:
            total = len(source_pdf)
            if total == 0:
                raise RuntimeError("PDF 无页面")
            if total > MAX_PAGES:
                raise RuntimeError(f"PDF 超过 {MAX_PAGES} 页上限")
            scale = dpi_clamped / 72.0

            output_pdf = pikepdf.Pdf.new()
            for index in range(total):
                _set_job(
                    job_id,
                    progress=10 + int(80 * (index + 1) / total),
                    stage=f"recompressing {index + 1}/{total}",
                )
                page = source_pdf[index]
                width_pt, height_pt = page.get_size()
                pil_image = page.render(scale=scale).to_pil().convert("RGB")

                # 量化降色：大幅减少 JPEG 编码复杂度（文字/图表页尤其受益）
                w, h = pil_image.size
                if w * h > 64000 and quality_clamped <= 80:
                    pil_image = pil_image.quantize(
                        colors=256, method=Image.Quantize.FASTOCTREE
                    ).convert("RGB")

                buf = _io.BytesIO()
                pil_image.save(
                    buf, format="JPEG", quality=quality_clamped, optimize=True, subsampling=0
                )
                jpeg_bytes = buf.getvalue()
                img_w, img_h = pil_image.size

                # DCTDecode 直嵌 JPEG 字节流——不经过任何解码/重编码
                img_stream = output_pdf.make_stream(
                    jpeg_bytes,
                    Dictionary(
                        Type=Name.XObject,
                        Subtype=Name.Image,
                        Width=img_w,
                        Height=img_h,
                        ColorSpace=Name.DeviceRGB,
                        BitsPerComponent=8,
                        Filter=Name.DCTDecode,
                    ),
                )
                # 内容流：把图像缩放到整页
                content = f"q {width_pt} 0 0 {height_pt} 0 0 cm /Im0 Do Q"
                content_stream = output_pdf.make_stream(content.encode("ascii"))

                page_dict = output_pdf.make_indirect(
                    Dictionary(
                        Type=Name.Page,
                        MediaBox=Array([0, 0, width_pt, height_pt]),
                        Contents=content_stream,
                        Resources=Dictionary(
                            XObject=Dictionary(Im0=img_stream),
                        ),
                    )
                )
                output_pdf.pages.append(pikepdf.Page(page_dict))
        finally:
            source_pdf.close()

        output_pdf.save(str(out_path))
        compressed_bytes = out_path.stat().st_size

        # 保底返原：压缩后体积反而更大时退回原文件（文字型 PDF 栅格化后矢量
        # 信息丢失，像素量可能超过原文件紧凑的内容流编码）
        if compressed_bytes >= original_bytes:
            out_path.unlink(missing_ok=True)
            shutil.copy2(source, out_path)
            compressed_bytes = original_bytes
            ratio = 100.0
            note = (
                f"压缩无效（原 {original_bytes // 1024} KB，栅格化后体积不降）。"
                "已退回原始文件——文字型 PDF 建议先转为更紧凑格式（如 PDF/A）再压缩。"
            )
            engine = "fallback-original"
        else:
            ratio = round(compressed_bytes / max(original_bytes, 1) * 100, 1)
            note = (
                f"栅格化深度压缩：视觉版式还原，无可编辑文本层"
                f"（原 {original_bytes // 1024} KB → {compressed_bytes // 1024} KB，{ratio}%）。"
                "适合扫描件/图片型 PDF；文字型 PDF 建议用前端轻压缩"
            )
            engine = "pypdfium2+pikepdf+DCTDecode"

        _set_job(
            job_id,
            status="done",
            progress=100,
            stage="done",
            engine=engine,
            page_count=total,
            dpi=dpi_clamped,
            quality=quality_clamped,
            original_bytes=original_bytes,
            compressed_bytes=compressed_bytes,
            compression_ratio_pct=ratio,
            note=note,
            outputs=[{"name": out_path.name, "dir": str(out_path.parent)}],
        )

    _spawn_job(job_id, tmp_dir, target)
    return {"job_id": job_id, "status": "running"}
    return {"job_id": job_id, "status": "running"}


# ---------------------------------------------------------------------------
# HTML / Markdown → PDF（reportlab platypus，HTML/markdown 子集解析，零联网）
# ---------------------------------------------------------------------------

_PDF_FONT_STATE: dict[str, Any] = {}


def _ensure_pdf_font() -> str:
    """注册可写的 CJK 字体：优先 Windows 雅黑（内嵌子集），否则 Adobe CID 兜底。"""
    if "font" in _PDF_FONT_STATE:
        return str(_PDF_FONT_STATE["font"])
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfbase.ttfonts import TTFont

    font = "Helvetica"
    for candidate in (r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\simhei.ttf"):
        try:
            if Path(candidate).exists():
                name = "MSyh" if candidate.endswith(".ttc") else "SimHei"
                if name not in pdfmetrics.getRegisteredFontNames():
                    pdfmetrics.registerFont(TTFont(name, candidate, subfontIndex=0))
                font = name
                break
        except Exception:  # noqa: BLE001
            continue
    if font == "Helvetica":
        try:
            if "STSong-Light" not in pdfmetrics.getRegisteredFontNames():
                pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            font = "STSong-Light"
        except Exception:  # noqa: BLE001
            font = "Helvetica"
    _PDF_FONT_STATE["font"] = font
    return font


def _esc(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _inline_markup(text: str) -> str:
    """内联标记 → reportlab Paragraph 标记（b/i/code/u）。"""
    out = _esc(text)
    out = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", out, flags=re.S)
    out = re.sub(r"(?<!\*)\*([^*\n]+)\*(?!\*)", r"<i>\1</i>", out)
    out = re.sub(r"`([^`]+)`", r'<font face="Courier" color="#7a1f1f">\1</font>', out)
    return out


class _HtmlSubsetParser(HTMLParser):
    """HTML 子集 → platypus flowables（h1-h6/p/div/br/ul/ol/li/table/blockquote/
    pre/hr/b/strong/i/em/code/u；img 忽略并记录提示；不联网取任何资源）。"""

    def __init__(self, font: str) -> None:
        super().__init__(convert_charrefs=True)
        from reportlab.lib import colors
        from reportlab.lib.styles import ParagraphStyle

        self.font = font
        self.warnings: list[str] = []
        self.flowables: list[Any] = []
        base = ParagraphStyle("base", fontName=font, fontSize=11, leading=17, textColor=colors.HexColor("#141414"))
        self.styles = {
            level: ParagraphStyle(
                f"h{level}", parent=base,
                fontSize=max(20 - level * 2, 12), leading=max(24 - level * 2, 15),
                spaceBefore=10 if level <= 1 else 6, spaceAfter=4,
            )
            for level in range(1, 7)
        }
        self.p_style = base
        self.quote_style = ParagraphStyle("quote", parent=base, leftIndent=16, textColor=colors.HexColor("#555555"))
        self.list_styles: dict[str, ParagraphStyle] = {}
        self._tag_stack: list[str] = []
        self._buf: list[str] = []
        self._list_ctx: list[dict[str, Any]] = []
        self._table_rows: list[list[str]] = []
        self._in_pre = False
        self._pre_buf: list[str] = []

    # -- helpers --
    def _flush_paragraph(self) -> None:
        text = "".join(self._buf).strip()
        self._buf = []
        if not text:
            return
        if self._table_rows:
            # 表格单元格里出现块级文本（少见）：按段落并入当前单元格
            self._table_rows[-1].append(text)
            return
        if self._list_ctx:
            ctx = self._list_ctx[-1]
            ctx["items"].append(text)
            return
        tag = self._tag_stack[-1] if self._tag_stack else "p"
        if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            self.flowables.append(Paragraph(text, self.styles[int(tag[1])]))
        elif tag == "blockquote":
            self.flowables.append(Paragraph(text, self.quote_style))
        else:
            self.flowables.append(Paragraph(text, self.p_style))

    def _flush_table(self) -> None:
        from reportlab.lib import colors
        from reportlab.platypus import Table

        rows = self._table_rows
        self._table_rows = []
        if not rows:
            return
        width = max(len(r) for r in rows)
        data = [r + [""] * (width - len(r)) for r in rows]
        table = Table(data, hAlign="LEFT")
        table.setStyle(
            [
                ("GRID", (0, 0), (-1, -1), 0.6, colors.HexColor("#999999")),
                ("FONTNAME", (0, 0), (-1, -1), self.font),
                ("FONTSIZE", (0, 0), (-1, -1), 9.5),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f0e6c8")),
            ]
        )
        self.flowables.append(table)
        from reportlab.platypus import Spacer

        self.flowables.append(Spacer(1, 8))

    def _flush_pre(self) -> None:
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.platypus import Preformatted

        text = "".join(self._pre_buf).strip("\n")
        self._pre_buf = []
        self._in_pre = False
        if text.strip():
            style = ParagraphStyle("pre", fontName="Courier", fontSize=9, leading=12)
            self.flowables.append(Preformatted(text, style))

    # -- HTMLParser 协议 --
    def handle_starttag(self, tag: str, attrs: Any) -> None:
        if tag in ("h1", "h2", "h3", "h4", "h5", "h6", "p", "div", "blockquote", "li"):
            self._flush_paragraph()
        if tag in ("ul", "ol"):
            self._list_ctx.append({"ordered": tag == "ol", "items": []})
        elif tag == "table":
            self._table_rows = []
        elif tag == "tr":
            self._flush_paragraph()
            self._table_rows.append([])
        elif tag == "pre":
            self._in_pre = True
        elif tag == "br":
            self._buf.append("<br/>")
        elif tag == "hr":
            from reportlab.platypus import HRFlowable

            self._flush_paragraph()
            self.flowables.append(HRFlowable(width="100%", thickness=0.8))
        elif tag == "img":
            self.warnings.append("忽略图片（离线文本排版不含图片）")
        if tag in ("b", "strong"):
            self._buf.append("<b>")
            self._tag_stack.append(tag)
        elif tag in ("i", "em"):
            self._buf.append("<i>")
            self._tag_stack.append(tag)
        elif tag == "code" and not self._in_pre:
            self._buf.append('<font face="Courier">')
            self._tag_stack.append(tag)
        elif tag == "u":
            self._buf.append("<u>")
            self._tag_stack.append(tag)
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6", "p", "div", "blockquote", "li"):
            self._tag_stack.append(tag)

    def handle_endtag(self, tag: str) -> None:
        if tag in ("b", "strong") and "<b>" in self._buf:
            self._buf.append("</b>")
            if tag in self._tag_stack:
                self._tag_stack.reverse()
                self._tag_stack.remove(tag)
                self._tag_stack.reverse()
        elif tag in ("i", "em") and "<i>" in self._buf:
            self._buf.append("</i>")
            if tag in self._tag_stack:
                self._tag_stack.reverse()
                self._tag_stack.remove(tag)
                self._tag_stack.reverse()
        elif tag == "code":
            self._buf.append("</font>")
            if "code" in self._tag_stack:
                self._tag_stack.reverse()
                self._tag_stack.remove("code")
                self._tag_stack.reverse()
        elif tag == "u":
            self._buf.append("</u>")
            if "u" in self._tag_stack:
                self._tag_stack.reverse()
                self._tag_stack.remove("u")
                self._tag_stack.reverse()
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6", "p", "div", "blockquote"):
            self._flush_paragraph()
            if tag in self._tag_stack:
                self._tag_stack.reverse()
                self._tag_stack.remove(tag)
                self._tag_stack.reverse()
        elif tag == "li":
            self._flush_paragraph()
            if "li" in self._tag_stack:
                self._tag_stack.reverse()
                self._tag_stack.remove("li")
                self._tag_stack.reverse()
        elif tag in ("ul", "ol"):
            self._flush_paragraph()
            if self._list_ctx:
                ctx = self._list_ctx.pop()
                for idx, item in enumerate(ctx["items"], start=1):
                    bullet = f"{idx}." if ctx["ordered"] else "•"
                    style = self.p_style
                    self.flowables.append(Paragraph(f"{bullet}&nbsp;&nbsp;{item}", style))
        elif tag == "table":
            self._flush_table()
        elif tag in ("td", "th"):
            # 单元格文本（_buf）直接落行，不经段落上下文
            text = "".join(self._buf).strip()
            self._buf = []
            if self._table_rows:
                self._table_rows[-1].append(text)
        elif tag == "pre":
            self._flush_pre()

    def handle_data(self, data: str) -> None:
        if self._in_pre:
            self._pre_buf.append(data)
            return
        text = _esc(data) if ("<" in data or ">" in data or "&" in data) else data
        self._buf.append(data)


def _html_to_flowables(html_text: str, font: str) -> tuple[list[Any], list[str]]:
    parser = _HtmlSubsetParser(font)
    # 给非完整 HTML 的纯文本一个最小包裹，避免纯文本走 HTML 解析丢换行
    stripped = html_text.strip()
    if "<" not in stripped:
        return _markdown_to_flowables(stripped, font)
    parser.feed(stripped)
    parser.close()
    parser._flush_paragraph()
    parser._flush_table()
    if parser._pre_buf:
        parser._flush_pre()
    return parser.flowables, parser.warnings


def _markdown_to_flowables(md_text: str, font: str) -> tuple[list[Any], list[str]]:
    """Markdown 子集（标题/列表/引用/代码块/分隔线/粗斜体内联）。"""
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import HRFlowable, Preformatted, Spacer

    flowables: list[Any] = []
    base = ParagraphStyle("md", fontName=font, fontSize=11, leading=17)
    code_style = ParagraphStyle("mdpre", fontName="Courier", fontSize=9, leading=12)
    quote_style = ParagraphStyle("mdq", parent=base, leftIndent=16, textColor=colors.HexColor("#555555"))
    lines = md_text.replace("\r\n", "\n").split("\n")
    para: list[str] = []
    in_code = False
    code_buf: list[str] = []

    def flush_para() -> None:
        if para:
            flowables.append(Paragraph(_inline_markup(" ".join(para)), base))
            para.clear()

    for line in lines:
        if line.strip().startswith("```"):
            if in_code:
                flowables.append(Preformatted("\n".join(code_buf), code_style))
                code_buf.clear()
                in_code = False
            else:
                flush_para()
                in_code = True
            continue
        if in_code:
            code_buf.append(line)
            continue
        stripped = line.strip()
        if not stripped:
            flush_para()
            continue
        m = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m:
            flush_para()
            level = len(m.group(1))
            flowables.append(
                Paragraph(
                    _inline_markup(m.group(2)),
                    ParagraphStyle(f"mdh{level}", parent=base, fontSize=max(20 - level * 2, 12), leading=max(24 - level * 2, 15), spaceBefore=8),
                )
            )
            continue
        if stripped in ("---", "***", "___"):
            flush_para()
            flowables.append(HRFlowable(width="100%", thickness=0.8))
            continue
        m = re.match(r"^>\s?(.*)$", stripped)
        if m:
            flush_para()
            flowables.append(Paragraph(_inline_markup(m.group(1)), quote_style))
            continue
        m = re.match(r"^[-*+]\s+(.*)$", stripped)
        if m:
            flush_para()
            flowables.append(Paragraph(f"•&nbsp;&nbsp;{_inline_markup(m.group(1))}", base))
            continue
        m = re.match(r"^\d+[.)]\s+(.*)$", stripped)
        if m:
            flush_para()
            num = re.match(r"^\d+", stripped).group(0)
            flowables.append(Paragraph(f"{num}.&nbsp;&nbsp;{_inline_markup(m.group(1))}", base))
            continue
        para.append(stripped)
    if in_code and code_buf:
        flowables.append(Preformatted("\n".join(code_buf), code_style))
    flush_para()
    return flowables, []


def _build_pdf_from_flowables(flowables: list[Any], out_path: Path, title: str) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Spacer

    doc = SimpleDocTemplate(
        str(out_path),
        pagesize=A4,
        title=title,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
    )
    doc.build(flowables or [Spacer(1, 1)])


@router.post("/html-to-pdf")
def html_to_pdf(
    file: Optional[UploadFile] = File(None),
    content: str = Form(""),
    title: str = Form("Document"),
) -> dict[str, Any]:
    """HTML/Markdown 子集 → PDF（同步端点；接受粘贴内容或 .html/.md/.txt 文件）。"""
    is_html = False
    if file is not None and (file.filename or "").strip():
        raw_name = Path(file.filename).name
        ext = Path(raw_name).suffix.lower()
        if ext not in (".html", ".htm", ".md", ".markdown", ".txt"):
            raise HTTPException(status_code=400, detail=f"不支持的格式: {ext or '(无扩展名)'}，支持 .html/.md/.txt 或直接粘贴内容")
        raw = file.file.read()
        if len(raw) > 4 * 1024 * 1024:
            raise HTTPException(status_code=400, detail="文本超过 4MB 上限")
        text = raw.decode("utf-8", errors="replace")
        is_html = ext in (".html", ".htm")
        base = _safe_base_name(raw_name)
    else:
        text = (content or "").strip()
        if not text:
            raise HTTPException(status_code=400, detail="请上传 .html/.md/.txt 文件或粘贴内容")
        is_html = bool(re.search(r"<(html|body|p|h[1-6]|div|ul|ol|table|br)\b", text, re.I))
        base = "pasted"

    out_path = _unique_output_path(f"{base}_pdf", ".pdf")
    font = _ensure_pdf_font()
    if is_html:
        flowables, warnings = _html_to_flowables(text, font)
    else:
        flowables, warnings = _markdown_to_flowables(text, font)
    if not flowables:
        raise HTTPException(status_code=400, detail="未解析到可排版内容")
    _build_pdf_from_flowables(flowables, out_path, title or base)
    return {
        "status": "success",
        "engine": "reportlab platypus",
        "input_format": "html" if is_html else "markdown",
        "warnings": warnings,
        "download_name": out_path.name,
        "output_dir": str(out_path.parent),
    }


# ---------------------------------------------------------------------------
# OCR 导出（RapidOCR → TXT / 可检索夹心 PDF：原页图 + 隐形文字层）
# ---------------------------------------------------------------------------

@router.post("/ocr-export")
def ocr_export(
    file: UploadFile = File(...),
    output: str = Form("txt"),
    dpi: int = Form(200),
) -> dict[str, Any]:
    fmt = "pdf" if output == "pdf" else "txt"
    dpi_clamped = min(max(int(dpi), 72), 300)

    tmp_dir = Path(tempfile.mkdtemp(prefix="convert_ocr_"))
    source, base = _save_pdf_upload(tmp_dir, file)
    suffix = "_ocr.txt" if fmt == "txt" else "_ocr.pdf"
    out_path = _unique_output_path(f"{base}{suffix}", "")
    job_id = _new_job("ocr-export")

    def target() -> None:
        import io as _io

        import numpy as np
        import pypdfium2 as pdfium
        from PIL import Image
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas as rl_canvas

        try:
            from core.detector.ocr_channel import get_ocr_engine

            engine = get_ocr_engine()
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError(f"OCR 引擎不可用: {exc}") from exc

        font = _ensure_pdf_font()
        scale = dpi_clamped / 72.0
        pages_text: list[list[str]] = []
        char_count = 0

        source_pdf = pdfium.PdfDocument(str(source))
        try:
            total = len(source_pdf)
            if total == 0:
                raise RuntimeError("PDF 无页面")
            if total > MAX_PAGES:
                raise RuntimeError(f"PDF 超过 {MAX_PAGES} 页上限")
            c = rl_canvas.Canvas(str(out_path)) if fmt == "pdf" else None
            for index in range(total):
                _set_job(job_id, progress=10 + int(80 * (index + 1) / total), stage=f"recognizing {index + 1}/{total}")
                page = source_pdf[index]
                width_pt, height_pt = page.get_size()
                arr = np.asarray(page.render(scale=scale).to_pil().convert("RGB"))
                result = engine(arr)
                lines: list[str] = []
                if result and result[0]:
                    for quad, text, _score in result[0]:
                        clean = (text or "").strip()
                        if clean:
                            lines.append(clean)
                            if c is not None:
                                xs = [p[0] for p in quad]
                                ys = [p[1] for p in quad]
                                x0, x1 = min(xs) / scale, max(xs) / scale
                                y_top, y_bot = min(ys) / scale, max(ys) / scale
                                box_h = max(y_bot - y_top, 4.0)
                                # 隐形文字层（渲染模式 3）：按行放在图像原位，可搜索可复制
                                t = c.beginText()
                                t.setTextRenderMode(3)
                                t.setFont(font, box_h * 0.85)
                                t.setTextOrigin(x0, height_pt - y_bot)
                                t.textOut(clean)
                                c.drawText(t)
                pages_text.append(lines)
                char_count += sum(len(l) for l in lines)
                if c is not None:
                    pil_image = Image.fromarray(arr).convert("RGB")
                    buf = _io.BytesIO()
                    pil_image.save(buf, format="JPEG", quality=80, optimize=True)
                    buf.seek(0)
                    c.setPageSize((width_pt, height_pt))
                    c.drawImage(ImageReader(buf), 0, 0, width=width_pt, height=height_pt)
                    c.showPage()
            if c is not None:
                c.save()
        finally:
            source_pdf.close()

        if fmt == "txt":
            body = "\n\n".join(
                f"---- Page {i + 1} ----\n" + "\n".join(lines or [""])
                for i, lines in enumerate(pages_text)
            )
            Path(out_path).write_text(body, encoding="utf-8")

        if char_count == 0:
            raise RuntimeError("未识别到任何文字（页面可能为空白或图像质量过低）")
        _set_job(
            job_id, status="done", progress=100, stage="done",
            engine="RapidOCR+pypdfium2" + ("+reportlab" if fmt == "pdf" else ""),
            page_count=total, char_count=char_count, dpi=dpi_clamped,
            note=(
                "夹心 PDF：原页视觉 + 隐形文字层（可搜索/复制），文字来自 OCR 识别，"
                "准确率受图像质量与语种影响" if fmt == "pdf" else "纯文本导出（OCR 识别结果）"
            ),
            outputs=[{"name": out_path.name, "dir": str(out_path.parent)}],
        )

    _spawn_job(job_id, tmp_dir, target)
    return {"job_id": job_id, "status": "running"}
